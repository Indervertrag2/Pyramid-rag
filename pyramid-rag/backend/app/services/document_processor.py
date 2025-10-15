"""
Advanced Document Processing Pipeline for RAG System
Implements SHA-256 deduplication, metadata extraction, text chunking, and embeddings.
"""

import hashlib
import logging
logger = logging.getLogger(__name__)
import os
import uuid
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple
from datetime import datetime
import mimetypes

# Document processing
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False

try:
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Language and metadata
try:
    from langdetect import detect, LangDetectError
    HAS_LANGDETECT = True
except ImportError:
    HAS_LANGDETECT = False

try:
    import magic
    HAS_MAGIC = True
except ImportError:
    HAS_MAGIC = False

# ML Libraries (for embeddings)
try:
    from sentence_transformers import SentenceTransformer
    HAS_EMBEDDINGS = True
except ImportError:
    HAS_EMBEDDINGS = False

# OCR (when available)
try:
    from surya import OCR
    HAS_SURYA = True
except ImportError:
    HAS_SURYA = False

from app.models import FileType, Document, DocumentChunk, ChatFile, ChatFileChunk
from app.schemas import FileScopeEnum


class DocumentProcessor:
    """Advanced document processing with RAG optimization."""

    def __init__(self):
        self.upload_dir = Path("data/uploads")
        self.upload_dir.mkdir(parents=True, exist_ok=True)

        # Initialize embedding model if available
        self.embedding_model = None
        # ✅ Upgraded to BGE-M3 (1024 dimensions, best multilingual performance)
        self.embedding_model_name = os.getenv('EMBEDDING_MODEL', 'BAAI/bge-m3')
        preferred_device = os.getenv('EMBEDDING_DEVICE')
        if preferred_device:
            self.embedding_device = preferred_device
        else:
            try:
                import torch  # type: ignore
                self.embedding_device = 'cuda' if torch.cuda.is_available() else 'cpu'
            except Exception:
                self.embedding_device = 'cpu'

        try:
            self.chunk_size_words = int(os.getenv('DOC_CHUNK_SIZE_WORDS', os.getenv('EMBEDDING_CHUNK_SIZE', '512')))
        except (TypeError, ValueError):
            self.chunk_size_words = 512

        try:
            self.chunk_overlap_words = int(os.getenv('DOC_CHUNK_OVERLAP_WORDS', os.getenv('EMBEDDING_CHUNK_OVERLAP', '50')))
        except (TypeError, ValueError):
            self.chunk_overlap_words = 50

        if HAS_EMBEDDINGS:
            try:
                logger.info('Loading embedding model: %s (this may take a moment for BGE-M3)...', self.embedding_model_name)
                self.embedding_model = SentenceTransformer(
                    self.embedding_model_name,
                    device=self.embedding_device,
                    trust_remote_code=True  # Required for BGE-M3
                )
                logger.info('✅ Embedding model loaded: %s on %s (dimensions: %d)',
                           self.embedding_model_name,
                           self.embedding_device,
                           self.embedding_model.get_sentence_embedding_dimension())
            except Exception as e:
                logger.warning('Could not load embedding model %s: %s', self.embedding_model_name, e)

        # Initialize OCR if available
        self.ocr_engine = None
        if HAS_SURYA:
            try:
                self.ocr_engine = OCR()
                print("Surya OCR initialized")
            except Exception as e:
                print(f"Could not initialize OCR: {e}")

    def _sanitize_text(self, text: str) -> str:
        """Strip characters (like NULL) that cannot be stored in Postgres TEXT."""
        if not text:
            return text
        return text.replace('\x00', ' ')

    def calculate_file_hash(self, file_path: Path) -> str:
        """Calculate SHA-256 hash of file for deduplication."""
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                sha256_hash.update(chunk)
        return sha256_hash.hexdigest()

    def detect_file_type(self, file_path: Path, original_filename: str) -> Tuple[FileType, str]:
        """Detect file type and MIME type."""

        # Get file extension
        ext = file_path.suffix.lower().lstrip('.')

        # Detect MIME type
        mime_type = mimetypes.guess_type(str(file_path))[0]
        if not mime_type and HAS_MAGIC:
            try:
                mime_type = magic.from_file(str(file_path), mime=True)
            except:
                mime_type = "application/octet-stream"

        # Map to FileType enum
        type_mapping = {
            'pdf': FileType.PDF,
            'doc': FileType.WORD, 'docx': FileType.WORD,
            'xls': FileType.EXCEL, 'xlsx': FileType.EXCEL,
            'ppt': FileType.POWERPOINT, 'pptx': FileType.POWERPOINT,
            'txt': FileType.TEXT, 'md': FileType.TEXT, 'rst': FileType.TEXT,
            'jpg': FileType.IMAGE, 'jpeg': FileType.IMAGE, 'png': FileType.IMAGE,
            'gif': FileType.IMAGE, 'bmp': FileType.IMAGE, 'tiff': FileType.IMAGE,
            'mp4': FileType.VIDEO, 'avi': FileType.VIDEO, 'mov': FileType.VIDEO,
            'mp3': FileType.AUDIO, 'wav': FileType.AUDIO, 'flac': FileType.AUDIO,
            'dwg': FileType.CAD, 'dxf': FileType.CAD, 'step': FileType.CAD
        }

        detected_type = type_mapping.get(ext, FileType.OTHER)
        return detected_type, mime_type or "application/octet-stream"

    def extract_text_content(self, file_path: Path, file_type: FileType) -> Tuple[str, Dict]:
        """Extract text content from various file formats."""

        content = ""
        metadata = {"extraction_method": "unknown", "success": False}

        try:
            if file_type == FileType.PDF:
                content, metadata = self._extract_pdf_text(file_path)
            elif file_type == FileType.WORD and HAS_DOCX:
                content, metadata = self._extract_docx_text(file_path)
            elif file_type == FileType.EXCEL and HAS_OPENPYXL:
                content, metadata = self._extract_xlsx_text(file_path)
            elif file_type == FileType.POWERPOINT and HAS_PPTX:
                content, metadata = self._extract_pptx_text(file_path)
            elif file_type == FileType.TEXT:
                content, metadata = self._extract_plain_text(file_path)
            else:
                # Fallback: try to read as text
                content, metadata = self._extract_plain_text(file_path)

        except Exception as e:
            metadata = {
                "extraction_method": "error",
                "success": False,
                "error": str(e)
            }

        return content, metadata

    def _extract_pdf_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Extract text from PDF using available libraries."""
        warnings: List[str] = []

        if HAS_PYMUPDF:
            try:
                doc = fitz.open(str(file_path))
                text_content = []

                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    if text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{text}")

                page_count = len(doc)
                doc.close()
                content = "\n\n".join(text_content).strip()

                if content:
                    metadata = {
                        "extraction_method": "pymupdf",
                        "success": True,
                        "pages": page_count,
                        "character_count": len(content)
                    }
                    return content, metadata

                warnings.append("pymupdf_extracted_empty_text")

            except Exception as e:
                warnings.append(f"pymupdf_error: {e}")

        if HAS_PYPDF:
            try:
                reader = PdfReader(str(file_path))
                text_content = []
                for page_idx, page in enumerate(reader.pages, 1):
                    try:
                        text = page.extract_text() or ""
                    except Exception as page_error:
                        warnings.append(f"pypdf_page_error_{page_idx}: {page_error}")
                        text = ""
                    text = text.strip()
                    if text:
                        text_content.append(f"[Page {page_idx}]\n{text}")

                content = "\n\n".join(text_content).strip()
                if content:
                    metadata = {
                        "extraction_method": "pypdf",
                        "success": True,
                        "pages": len(reader.pages),
                        "character_count": len(content)
                    }
                    if warnings:
                        metadata["warnings"] = warnings
                    return content, metadata

                warnings.append("pypdf_extracted_empty_text")

            except Exception as e:
                warnings.append(f"pypdf_error: {e}")

        content, metadata = self._extract_plain_text(file_path)
        if warnings:
            metadata = dict(metadata)
            metadata["warnings"] = warnings
        return content, metadata

    def _extract_docx_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Extract text from DOCX files."""
        try:
            doc = DocxDocument(str(file_path))
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            content = "\n".join(paragraphs)

            metadata = {
                "extraction_method": "python-docx",
                "success": True,
                "paragraphs": len(paragraphs),
                "character_count": len(content)
            }

            return content, metadata

        except Exception as e:
            return "", {
                "extraction_method": "docx_error",
                "success": False,
                "error": str(e)
            }

    def _extract_xlsx_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Extract text from Excel files."""
        try:
            workbook = openpyxl.load_workbook(str(file_path), data_only=True)
            content_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"[Sheet: {sheet_name}]")

                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        content_parts.append(row_text)

            content = "\n".join(content_parts)

            metadata = {
                "extraction_method": "openpyxl",
                "success": True,
                "sheets": len(workbook.sheetnames),
                "character_count": len(content)
            }

            return content, metadata

        except Exception as e:
            return "", {
                "extraction_method": "xlsx_error",
                "success": False,
                "error": str(e)
            }

    def _extract_pptx_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Extract text from PowerPoint files."""
        try:
            presentation = Presentation(str(file_path))
            content_parts = []

            for i, slide in enumerate(presentation.slides, 1):
                content_parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

            content = "\n".join(content_parts)

            metadata = {
                "extraction_method": "python-pptx",
                "success": True,
                "slides": len(presentation.slides),
                "character_count": len(content)
            }

            return content, metadata

        except Exception as e:
            return "", {
                "extraction_method": "pptx_error",
                "success": False,
                "error": str(e)
            }

    def _extract_plain_text(self, file_path: Path) -> Tuple[str, Dict]:
        """Extract plain text content."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()

                    metadata = {
                        "extraction_method": f"text_{encoding}",
                        "success": True,
                        "encoding": encoding,
                        "character_count": len(content)
                    }

                    return content, metadata

                except UnicodeDecodeError:
                    continue

            # If all encodings fail
            return "", {
                "extraction_method": "text_error",
                "success": False,
                "error": "Could not decode text with any encoding"
            }

        except Exception as e:
            return "", {
                "extraction_method": "text_error",
                "success": False,
                "error": str(e)
            }

    def detect_language(self, text: str) -> str:
        """Detect text language."""
        if not HAS_LANGDETECT or not text.strip():
            return "unknown"

        try:
            # Take a sample if text is very long
            sample_text = text[:1000] if len(text) > 1000 else text
            detected_lang = detect(sample_text)
            return detected_lang
        except LangDetectError:
            return "unknown"

    def extract_metadata(self, file_path: Path, content: str) -> Dict[str, Any]:
        """Extract comprehensive metadata from file and content."""

        file_stats = file_path.stat()

        metadata = {
            "file_size": file_stats.st_size,
            "created_time": datetime.fromtimestamp(file_stats.st_ctime).isoformat(),
            "modified_time": datetime.fromtimestamp(file_stats.st_mtime).isoformat(),
            "character_count": len(content),
            "word_count": len(content.split()) if content else 0,
            "language": self.detect_language(content),
            "has_content": bool(content.strip())
        }

        # Add file-specific metadata
        if HAS_MAGIC:
            try:
                file_info = magic.from_file(str(file_path))
                metadata["file_description"] = file_info
            except:
                pass

        return metadata

    def chunk_text(self, text: str, chunk_size: Optional[int] = None, overlap: Optional[int] = None) -> List[Dict]:
        """Intelligent text chunking for RAG."""
        if not text.strip():
            return []

        effective_chunk = chunk_size or self.chunk_size_words
        effective_overlap = overlap or self.chunk_overlap_words
        effective_overlap = min(effective_overlap, effective_chunk - 1) if effective_chunk > 1 else 0

        chunks: List[Dict[str, Any]] = []
        words = text.split()
        step = max(1, effective_chunk - effective_overlap)

        for i in range(0, len(words), step):
            chunk_words = words[i:i + effective_chunk]
            chunk_text = " ".join(chunk_words)
            chunk_text = self._sanitize_text(chunk_text)

            if chunk_text.strip():
                chunk_info = {
                    "content": chunk_text,
                    "start_word": i,
                    "end_word": i + len(chunk_words),
                    "word_count": len(chunk_words),
                    "character_count": len(chunk_text)
                }
                chunks.append(chunk_info)

        return chunks

    def generate_embeddings(self, text_chunks: List[str]) -> List[List[float]]:
        """Generate embeddings for text chunks."""
        if not self.embedding_model or not text_chunks:
            return []

        try:
            embeddings = self.embedding_model.encode(text_chunks)
            return embeddings.tolist()
        except Exception as e:
            print(f"âŒ Embedding generation failed: {e}")
            return []

    async def process_document(
        self,
        file_path: Path,
        original_filename: str,
        scope: FileScopeEnum = FileScopeEnum.GLOBAL,
        generate_embeddings: bool = True
    ) -> Dict[str, Any]:
        """
        Complete document processing pipeline.

        Returns processing results with extracted content, metadata, chunks, and embeddings.
        """

        result = {
            "success": False,
            "file_hash": None,
            "file_type": None,
            "mime_type": None,
            "content": "",
            "language": "unknown",
            "metadata": {},
            "chunks": [],
            "embeddings": [],
            "processing_time": 0,
            "errors": []
        }

        start_time = datetime.now()

        try:
            # 1. Calculate file hash for deduplication
            result["file_hash"] = self.calculate_file_hash(file_path)

            # 2. Detect file type and MIME type
            file_type, mime_type = self.detect_file_type(file_path, original_filename)
            result["file_type"] = file_type
            result["mime_type"] = mime_type

            # 3. Extract text content
            content, extraction_metadata = self.extract_text_content(file_path, file_type)
            content = self._sanitize_text(content)
            result["content"] = content

            # 4. Detect language
            result["language"] = self.detect_language(content)

            # 5. Extract comprehensive metadata
            file_metadata = self.extract_metadata(file_path, content)
            result["metadata"] = {**file_metadata, **extraction_metadata}

            # 6. Generate text chunks
            if content.strip():
                chunks = self.chunk_text(content)
                result["chunks"] = chunks

                # 7. Generate embeddings (only if requested)
                if generate_embeddings and chunks and self.embedding_model:
                    chunk_texts = [chunk["content"] for chunk in chunks]
                    embeddings = self.generate_embeddings(chunk_texts)
                    result["embeddings"] = embeddings
                    if embeddings and self.embedding_model_name:
                        metadata = result.setdefault("metadata", {})
                        metadata.setdefault("embedding_model", self.embedding_model_name)

            # 8. Calculate processing time
            end_time = datetime.now()
            result["processing_time"] = (end_time - start_time).total_seconds()

            result["success"] = True

        except Exception as e:
            result["errors"].append(f"Processing error: {str(e)}")
            print(f"âŒ Document processing failed: {e}")

        return result


# Global instance
document_processor = DocumentProcessor()



