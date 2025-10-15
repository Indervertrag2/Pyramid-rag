import os
import hashlib
from typing import Optional, List, Dict, Any
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import chardet
import json
from datetime import datetime
from app.utils.file_security import sanitize_filename, secure_join
from sqlalchemy.orm import Session

from app.models import Document, DocumentChunk, FileType
from app.database import get_db
try:
    from app.embeddings_service import embeddings_service
    EMBEDDINGS_AVAILABLE = True
except Exception as e:
    print(f"Warning: Embeddings service not available: {e}")
    embeddings_service = None
    EMBEDDINGS_AVAILABLE = False

class DocumentProcessor:
    """Process various document types and extract content"""

    def __init__(self):
        self.upload_dir = Path("data/uploads")
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.chunk_size = 1000  # Characters per chunk
        self.chunk_overlap = 200  # Overlap between chunks

    async def save_file(self, file) -> str:
        """Save uploaded file to disk"""
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        safe_original = sanitize_filename(file.filename or "upload")
        filename = f"{timestamp}_{safe_original}"
        file_path = secure_join(self.upload_dir, filename, fallback_prefix="upload")

        content = await file.read()
        with open(file_path, "wb") as f:
            f.write(content)

        return str(file_path)

    def get_file_type(self, filename: str) -> FileType:
        """Determine file type from extension"""
        ext = filename.lower().split('.')[-1]

        type_mapping = {
            'pdf': FileType.PDF,
            'doc': FileType.WORD,
            'docx': FileType.WORD,
            'xls': FileType.EXCEL,
            'xlsx': FileType.EXCEL,
            'ppt': FileType.POWERPOINT,
            'pptx': FileType.POWERPOINT,
            'txt': FileType.TEXT,
            'md': FileType.TEXT,
            'csv': FileType.TEXT,
            'jpg': FileType.IMAGE,
            'jpeg': FileType.IMAGE,
            'png': FileType.IMAGE,
            'gif': FileType.IMAGE,
            'bmp': FileType.IMAGE,
            'dwg': FileType.CAD,
            'dxf': FileType.CAD,
            'step': FileType.CAD,
            'stp': FileType.CAD,
            'iges': FileType.CAD,
            'mp4': FileType.VIDEO,
            'avi': FileType.VIDEO,
            'mkv': FileType.VIDEO,
            'mov': FileType.VIDEO,
            'mp3': FileType.AUDIO,
            'wav': FileType.AUDIO,
            'flac': FileType.AUDIO
        }

        return type_mapping.get(ext, FileType.OTHER)

    def extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file, with OCR fallback for scanned PDFs"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    text += page_text + "\n"

            # Check if we got meaningful text
            if len(text.strip()) < 50:  # Less than 50 chars probably means scanned PDF
                print(f"PDF appears to be scanned, attempting OCR...")
                try:
                    import pytesseract
                    from pdf2image import convert_from_path

                    # Convert PDF pages to images
                    images = convert_from_path(file_path, dpi=200)
                    ocr_text = ""

                    for i, image in enumerate(images):
                        # OCR with German + English
                        page_text = pytesseract.image_to_string(image, lang='deu+eng')
                        ocr_text += f"\n--- Page {i+1} ---\n{page_text}"

                    if ocr_text.strip():
                        return f"[OCR Extracted from PDF]\n{ocr_text}"
                except Exception as ocr_e:
                    print(f"OCR fallback failed: {ocr_e}")

        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return ""
        return text

    def extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = DocxDocument(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except Exception as e:
            print(f"Error extracting DOCX text: {e}")
            return ""

    def extract_text_from_xlsx(self, file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            workbook = load_workbook(file_path, data_only=True)
            text = ""
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\nSheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join([str(cell) if cell else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
            return text
        except Exception as e:
            print(f"Error extracting XLSX text: {e}")
            return ""

    def extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            presentation = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"\nSlide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error extracting PPTX text: {e}")
            return ""

    def extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            # Detect encoding
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'

            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text file: {e}")
            return ""

    def extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image file using OCR (Tesseract)"""
        try:
            import pytesseract
            from PIL import Image

            # Open the image
            img = Image.open(file_path)

            # Try OCR with German and English
            try:
                # Configure for German + English
                text = pytesseract.image_to_string(img, lang='deu+eng')
                if text.strip():
                    return f"[OCR Extracted Text]\n{text}"
            except Exception as ocr_error:
                print(f"OCR failed: {ocr_error}")

            # If OCR fails or returns nothing, return metadata
            return f"Image: {img.format}, Size: {img.size}, Mode: {img.mode}"
        except Exception as e:
            print(f"Error processing image: {e}")
            return ""

    def extract_text(self, file_path: str, file_type: FileType) -> str:
        """Extract text from any supported file type"""
        extractors = {
            FileType.PDF: self.extract_text_from_pdf,
            FileType.WORD: self.extract_text_from_docx,
            FileType.EXCEL: self.extract_text_from_xlsx,
            FileType.POWERPOINT: self.extract_text_from_pptx,
            FileType.TEXT: self.extract_text_from_txt,
            FileType.IMAGE: self.extract_text_from_image,
        }

        extractor = extractors.get(file_type)
        if extractor:
            return extractor(file_path)
        else:
            return f"File type {file_type.value} is not yet supported for text extraction"

    def chunk_text(self, text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []

        chunk_size = chunk_size or self.chunk_size
        overlap = overlap or self.chunk_overlap

        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = min(start + chunk_size, text_len)

            # Try to find a sentence boundary
            if end < text_len:
                for separator in ['. ', '! ', '? ', '\n\n', '\n', ' ']:
                    sep_pos = text.rfind(separator, start, end)
                    if sep_pos != -1:
                        end = sep_pos + len(separator)
                        break

            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - overlap if end < text_len else end

        return chunks

    async def process_document(self, document_id: str, db: Session):
        """Process a document: extract text and create chunks"""
        document = db.query(Document).filter(Document.id == document_id).first()
        if not document:
            return

        try:
            # Extract text
            text = self.extract_text(document.file_path, document.file_type)
            document.content = text[:10000]  # Store first 10k chars in document

            # Create chunks
            chunks = self.chunk_text(text)
            for idx, chunk_text in enumerate(chunks):
                chunk = DocumentChunk(
                    document_id=document.id,
                    chunk_index=idx,
                    content=chunk_text,
                    token_count=len(chunk_text.split()),
                    meta_data={
                        "page": idx // 3,  # Approximate page number
                        "total_chunks": len(chunks)
                    }
                )
                db.add(chunk)

            # Generate embeddings for the document chunks if available
            if EMBEDDINGS_AVAILABLE and embeddings_service:
                try:
                    print(f"Generating embeddings for document {document_id}...")
                    embeddings_success = await embeddings_service.process_document_embeddings(str(document.id), db)

                    if embeddings_success:
                        print(f"Embeddings generated successfully for document {document_id}")
                        document.processed = True
                    else:
                        print(f"Failed to generate embeddings for document {document_id}")
                        document.processing_error = "Failed to generate embeddings"
                        document.processed = True  # Mark as processed anyway
                except Exception as e:
                    print(f"Warning: Could not generate embeddings: {e}")
                    document.processed = True  # Mark as processed without embeddings
            else:
                print("Embeddings service not available, skipping embedding generation")
                document.processed = True

            db.commit()

        except Exception as e:
            document.processing_error = str(e)
            document.processed = False
            db.commit()
            print(f"Error processing document {document_id}: {e}")

    async def process_document_async(self, document_id: str):
        """Process document asynchronously (placeholder for Celery task)"""
        # In production, this would be a Celery task
        from app.database import SessionLocal
        db = SessionLocal()
        try:
            await self.process_document(document_id, db)
        finally:
            db.close()

    def get_document_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from document"""
        stat = os.stat(file_path)

        # Calculate file hash
        with open(file_path, 'rb') as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()

        return {
            "file_size": stat.st_size,
            "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "file_hash": file_hash
        }